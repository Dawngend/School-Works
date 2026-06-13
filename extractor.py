import os
import io
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Tesseract executable path for Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_pdf(file_path):
    extracted_text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                
                if text:
                    extracted_text += text + "\n\n"
                
                # Scan full-page screenshots
                if not text or len(text.strip()) < 50:
                    print(f"  -> Scanning image on PDF page {i + 1}...")
                    pil_image = page.to_image(resolution=300).original
                    ocr_text = pytesseract.image_to_string(pil_image)
                    
                    if ocr_text.strip():
                        extracted_text += "\n[Scanned from Image]:\n" + ocr_text + "\n\n"
                        
        return extracted_text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_pptx(file_path):
    extracted_text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
                
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"  -> Scanning image on PPTX slide {i + 1}...")
                    image_bytes = shape.image.blob
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(image)
                    
                    if ocr_text.strip():
                        extracted_text += "\n[Scanned from Image]:\n" + ocr_text + "\n"
                        
            extracted_text += "\n--- Next Slide ---\n\n"
        return extracted_text
    except Exception as e:
        return f"Error reading PPTX: {e}"

def process_module_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        print(f"Processing PDF: {file_path}...")
        return extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        print(f"Processing PPTX: {file_path}...")
        return extract_text_from_pptx(file_path)
    else:
        return "Unsupported file type. Please upload a PDF or PPTX."